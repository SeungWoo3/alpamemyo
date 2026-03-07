#!/usr/bin/env python3
"""
AVEES 템플릿 기반 Alpamayo-R1-10B 아키텍처 PPT 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy
from lxml import etree

TEMPLATE_PATH = "/home/seungwoo/workspace/docs/seminar/AVEES_Template.pptx"
OUTPUT_PATH = "/home/seungwoo/workspace/docs/seminar/architecture-overview.pptx"

# 색상 상수
COLOR_HEADER_BG = RGBColor(0x1B, 0x2A, 0x4A)
COLOR_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ROW_ODD = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ROW_EVEN = RGBColor(0xF0, 0xF0, 0xF0)
COLOR_BORDER = RGBColor(0xCC, 0xCC, 0xCC)


def set_cell_bg(cell, color):
    """셀 배경색 설정"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # 기존 solidFill 제거
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', str(color))


def set_cell_border(cell, color=COLOR_BORDER, width=12700):
    """셀 테두리 설정"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        existing = tcPr.findall(qn(edge))
        for e in existing:
            tcPr.remove(e)
        ln = etree.SubElement(tcPr, qn(edge))
        ln.set('w', str(width))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        solidFill = etree.SubElement(ln, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', str(color))


def style_table_header(table, col_count, font_size=12):
    """헤더 행 스타일 적용"""
    for col in range(col_count):
        cell = table.cell(0, col)
        set_cell_bg(cell, COLOR_HEADER_BG)
        set_cell_border(cell)
        tf = cell.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = COLOR_HEADER_FG
                run.font.size = Pt(font_size)


def style_table_data_rows(table, row_count, col_count, font_size=12):
    """데이터 행 교대 색상 적용"""
    for row in range(1, row_count):
        color = COLOR_ROW_ODD if row % 2 == 1 else COLOR_ROW_EVEN
        for col in range(col_count):
            cell = table.cell(row, col)
            set_cell_bg(cell, color)
            set_cell_border(cell)
            tf = cell.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(font_size)


def add_bold_label_text(ph, sections, font_size=14):
    """
    섹션별 라벨(볼드)+내용 텍스트 추가.
    sections: list of dict with 'label'(optional bold) and 'lines'(list of str/dict)
    **text** 구문으로 인라인 볼드 처리
    """
    tf = ph.text_frame
    tf.word_wrap = True
    tf.clear()

    first_para = True

    def get_para():
        nonlocal first_para
        if first_para:
            first_para = False
            return tf.paragraphs[0]
        return tf.add_paragraph()

    def add_inline_text(para, text, base_bold=False, fsize=font_size):
        """**..** 구문 파싱하여 인라인 볼드 처리"""
        if '**' in text:
            parts = text.split('**')
            is_bold_part = False
            for part in parts:
                if part:
                    run = para.add_run()
                    run.text = part
                    run.font.bold = base_bold or is_bold_part
                    run.font.size = Pt(fsize)
                is_bold_part = not is_bold_part
        else:
            run = para.add_run()
            run.text = text
            run.font.bold = base_bold
            run.font.size = Pt(fsize)

    for section in sections:
        label = section.get('label', '')
        lines = section.get('lines', [])
        after_space = section.get('after_space', True)

        if label:
            para = get_para()
            add_inline_text(para, label, base_bold=True)

        for line in lines:
            if isinstance(line, dict):
                text = line.get('text', '')
                bold = line.get('bold', False)
            else:
                text = line
                bold = False

            para = get_para()
            if not text:
                run = para.add_run()
                run.text = ''
                run.font.size = Pt(font_size)
            else:
                add_inline_text(para, text, base_bold=bold)

        if after_space:
            para = get_para()
            run = para.add_run()
            run.text = ''
            run.font.size = Pt(font_size)


def main():
    prs = Presentation(TEMPLATE_PATH)

    layouts = {layout.name: layout for layout in prs.slide_layouts}

    # ── 슬라이드 1: 표지 ──────────────────────────────────────────────
    slide = prs.slides.add_slide(layouts['표지'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    phs[10].text = "AVEES Lab"
    phs[0].text = "Alpamayo-R1-10B\n추론 파이프라인 아키텍처"
    phs[1].text = "7-Stage Module 상세 분석\n기능 · 입출력 차원 · 값의 의미"
    # 날짜는 텍스트박스로 추가 (placeholder idx=2가 슬라이드 인스턴스에 없음)
    from pptx.util import Emu as _Emu
    date_box = slide.shapes.add_textbox(_Emu(9448800), _Emu(6292875), _Emu(2743200), _Emu(465125))
    date_tf = date_box.text_frame
    date_run = date_tf.paragraphs[0].add_run()
    date_run.text = "2026-03-06"
    date_run.font.size = Pt(12)

    # ── 슬라이드 2: 7-Stage 파이프라인 요약 (테이블) ──────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "7-Stage 추론 파이프라인 요약"

    ph1 = phs[1]
    left = ph1.left
    top = ph1.top
    width = ph1.width
    height = ph1.height

    rows, cols = 8, 5
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    col_widths = [Emu(700000), Emu(1900000), Emu(3100000), Emu(2800000), Emu(3025958)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ['Stage', '모듈', '기능', '입력', '출력']
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h

    data = [
        ['1', 'Patch Embedding', '이미지→패치 임베딩', '(B,3,2,H,W)', '(N, 1152)'],
        ['2', 'Vision Transformer', '시각 특징 추출 (27L)', '(N, 1152)', '(N, 1152)'],
        ['3', 'Patch Merger', '공간 압축 + 차원 정합', '(N, 1152)', '(N/4, 4096)'],
        ['4', 'VLM (Qwen3-VL-8B)', 'CoT 추론 + KV cache', '(B, L, 4096)', 'KV cache (36L)'],
        ['5', 'Diffusion Head', '노이즈→action 생성', '(B,64,2) noise', '(B,64,2) denoised'],
        ['6', 'Action→Trajectory', 'Unicycle 적분', '(B,64,2)', '(B,64,3)'],
        ['7', 'Final Output', '다중 샘플 취합', '6× trajectory', '(B,1,6,64,3)'],
    ]
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val

    style_table_header(table, cols, font_size=12)
    style_table_data_rows(table, rows, cols, font_size=12)

    # placeholder 제거
    sp = ph1._element
    sp.getparent().remove(sp)

    # ── 슬라이드 3: Stage 1 Patch Embedding ──────────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "Stage 1: Patch Embedding (~10M params)"

    sections = [
        {
            'label': '기능',
            'lines': [
                '카메라 이미지를 16×16 비겹침 패치로 분할, Conv3d로 시간축 2프레임 동시 처리하여 1152차원 임베딩 생성',
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '입력: **(B×(H/16)×(W/16), 3, 2, 16, 16)**'},
                {'text': '· 추론 시 T=2 고정 (연속 2프레임 입력)'},
                {'text': '· RGB 3채널, temporal 2프레임, spatial 16×16 픽셀'},
                {'text': '· **값의 의미:** 정규화된 RGB 픽셀값 [0,1]'},
            ],
        },
        {
            'label': '출력 (두 가지)',
            'lines': [
                {'text': '① hidden_states: **(N_patches, 1152)** — Conv3d 출력 + 절대 위치 임베딩'},
                {'text': '   · N_patches = (H/16) × (W/16)'},
                {'text': '   · **값의 의미:** 시공간 패턴의 활성화 값 + 전역 위치 정보'},
                {'text': '② rotary_pos_emb: **(N_patches, 72)** — 패치 좌표(t,h,w)에서 계산된 RoPE 벡터'},
                {'text': '   · **값의 의미:** Attention Q/K에 적용할 상대 위치 회전 값 (매 레이어 사용)'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '**핵심 구조:** Conv3d(3→1152, kernel=2×16×16, stride=2×16×16) + 3D RoPE'},
            ],
            'after_space': False,
        },
    ]
    add_bold_label_text(phs[1], sections, font_size=15)

    # ── 슬라이드 4: Stage 2 Vision Transformer ───────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "Stage 2: Vision Transformer (27 Blocks, ~400M)"

    sections = [
        {
            'label': '기능',
            'lines': [
                '패치 임베딩 간 전역 관계 모델링. 16-head Attention으로 시각 패턴 포착, MLP로 비선형 변환',
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '입력: **(N_patches, 1152)** + RoPE 위치 인코딩'},
                {'text': '· **값의 의미:** 각 패치의 로컬 시각 특징 (Stage 1 출력)'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '출력: **(N_patches, 1152)**'},
                {'text': '· **값의 의미:** 이미지 전체 관계 (객체 간 상호작용, 공간 배치)가 인코딩된 글로벌 특징'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '**블록 구조:** LayerNorm → MHA(16 heads, head_dim=72) → LayerNorm → MLP(1152→4304→1152, SiLU)'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '**DeepStack:** 지정 중간 레이어에서 특징 추출 → VLM early layer에 주입'},
            ],
            'after_space': False,
        },
    ]
    add_bold_label_text(phs[1], sections, font_size=15)

    # ── 슬라이드 5: Stage 3 Patch Merger ─────────────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "Stage 3: Patch Merger (~37M params)"

    sections = [
        {
            'label': '기능',
            'lines': [
                '인접 2×2 패치를 하나로 병합 (시퀀스 1/4 축소) + VLM 입력 차원(4096)으로 투영',
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '입력: **(N_patches, 1152)**'},
                {'text': '· **값의 의미:** 전역 컨텍스트 반영된 개별 패치의 시각 특징'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '출력: **(N_patches/4, 4096)**'},
                {'text': '· **값의 의미:** 4개 패치를 통합한 고밀도 시각 표현, VLM hidden_size(4096)와 정합'},
                {'text': '· 예: 392 patches → 98 visual tokens'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '**구조:** 2×2 concat(1152×4=4608) → Linear(4608→4608) → GELU → Linear(4608→4096)'},
            ],
            'after_space': False,
        },
    ]
    add_bold_label_text(phs[1], sections, font_size=15)

    # ── 슬라이드 6: Stage 4-1 VLM 입력 구성 ─────────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "Stage 4-1: VLM 입력 시퀀스 구성"

    sections = [
        {
            'label': '기능',
            'lines': [
                '시각 토큰 + 텍스트 토큰 + 히스토리 토큰을 하나의 멀티모달 시퀀스로 결합',
            ],
        },
        {
            'label': '입력 구성',
            'lines': [
                {'text': '· Vision tokens: **(N/4, 4096)** — Stage 3 출력, 시각 정보'},
                {'text': '· Text tokens: **(N_text, 4096)** — 주행 명령어 임베딩 (예: "좌회전", "직진")'},
                {'text': '· Ego-history tokens: **(N_hist, 4096)** — fuse_traj_tokens으로 과거 waypoint 삽입'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '출력: **(B, L, 4096)** — L = N_vis + N_text + N_hist'},
                {'text': '· **값의 의미:** 멀티모달 컨텍스트의 통합 표현, VLM Transformer 입력'},
            ],
            'after_space': False,
        },
    ]
    add_bold_label_text(phs[1], sections, font_size=15)

    # ── 슬라이드 7: Stage 4-2 VLM Transformer ───────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "Stage 4-2: VLM Transformer (36L, ~8B) — 핵심 병목"

    sections = [
        {
            'label': '기능',
            'lines': [
                'GQA 기반 36-layer Transformer로 멀티모달 추론 수행 (Chain-of-Thought 생성)',
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '입력: **(B, L, 4096)** — 시각+언어+이력 통합 시퀀스'},
                {'text': '· **값의 의미:** 멀티모달 토큰 시퀀스'},
            ],
        },
        {
            'label': '출력',
            'lines': [
                {'text': '· Hidden states: **(B, L, 4096)** — 추론 결과'},
                {'text': '· KV cache: 36 layers × **(B, 8, L, 128)** — Expert Decoder cross-attention에 전달'},
                {'text': '· **값의 의미:** VLM reasoning context'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '**구조:** RMSNorm → GQA(32Q/8KV, head_dim=128) → RMSNorm → SwiGLU FFN(4096→12288→4096)'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '⚠ 전체 파라미터의 92%, 추론 시간의 78.5% — 메모리 최적화의 핵심 대상'},
                {'text': 'Autoregressive 생성: 토큰당 ~1.07s'},
            ],
            'after_space': False,
        },
    ]
    add_bold_label_text(phs[1], sections, font_size=14)

    # ── 슬라이드 8: Stage 5-1 Expert Decoder ────────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "Stage 5-1: Expert Decoder (16L, ~2B params)"

    sections = [
        {
            'label': '기능',
            'lines': [
                'VLM의 KV cache를 cross-attention으로 참조하면서 action 표현 정제',
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '입력: **(B, 64, 2048)** — action_in_proj 출력'},
                {'text': '· **값의 의미:** Fourier 인코딩된 timestep + noisy action의 고차원 표현'},
                {'text': '· + VLM KV cache (cross-attention context)'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '출력: **(B, 64, 2048)**'},
                {'text': '· **값의 의미:** VLM reasoning context를 반영한 64개 waypoint별 행동 표현'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '**구조:** RMSNorm → Self-Attn(16h, 2048) → Cross-Attn(VLM KV) → SwiGLU FFN(2048→8256→2048)'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '**특징:** non-causal attention (양방향), Diffusion의 10 Euler step마다 반복 호출'},
            ],
            'after_space': False,
        },
    ]
    add_bold_label_text(phs[1], sections, font_size=14)

    # ── 슬라이드 9: Stage 5-2 Diffusion Head ────────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "Stage 5-2: Diffusion Head (Flow Matching)"

    sections = [
        {
            'label': '기능',
            'lines': [
                '가우시안 노이즈에서 10 Euler step으로 action trajectory 생성',
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '입력: x₀ ~ N(0,I), shape **(B, 64, 2)**'},
                {'text': '· 64 waypoints × 2 (acceleration, curvature)'},
                {'text': '· **값의 의미:** 순수 랜덤 노이즈 (t=0)'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '출력: **(B, 64, 2)** — denoised action'},
                {'text': '· **값의 의미:** 정규화된 (acceleration, curvature), 물리 단위 복원 전'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '**Euler step:** x_{t+dt} = x_t + dt × v_θ(x_t, t), dt=0.1, 10 steps'},
            ],
        },
        {
            'label': 'step별 파이프라인:',
            'lines': [
                {'text': 'action_in_proj (Fourier+MLP → **(B,64,2048)**) → Expert Decoder (16L) → action_out_proj (Linear 2048→2)'},
            ],
            'after_space': False,
        },
    ]
    add_bold_label_text(phs[1], sections, font_size=14)

    # ── 슬라이드 10: Stage 6 Action→Trajectory ──────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "Stage 6: Action → Trajectory (Unicycle Kinematics)"

    sections = [
        {
            'label': '기능',
            'lines': [
                '(acceleration, curvature) action을 물리적으로 실현 가능한 경로로 변환',
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '입력: **(B, 64, 2)** — denoised action (정규화 상태)'},
                {'text': '· [:,:,0]: acceleration, [:,:,1]: curvature κ'},
                {'text': '· **값의 의미:** 차량의 가감속과 조향 곡률 (정규화 값)'},
            ],
        },
        {
            'label': '출력',
            'lines': [
                {'text': '· traj_xyz: **(B, 64, 3)** — [x(m), y(m), z(m)]'},
                {'text': '· traj_rot: **(B, 64, 3, 3)** — SO(3) 회전 행렬'},
                {'text': '· **값의 의미:** 차량 중심 로컬 좌표계, 6.4초 미래 경로 (dt=0.1s)'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '**변환 수식:** 역정규화 → v₀ 추정 → 속도 적분 → 방향각 적분(κ×v×dt) → 사다리꼴 위치 적분'},
            ],
            'after_space': False,
        },
    ]
    add_bold_label_text(phs[1], sections, font_size=14)

    # ── 슬라이드 11: Stage 7 Final Output ───────────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "Stage 7: Final Output (Multi-Sample)"

    sections = [
        {
            'label': '기능',
            'lines': [
                '다중 trajectory 샘플을 취합하여 최종 예측 결과 제공',
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '입력: 6× Stage 6 출력 (독립 노이즈에서 생성된 6개 후보)'},
            ],
        },
        {
            'label': '출력',
            'lines': [
                {'text': '· pred_xyz: **(B, 1, 6, 64, 3)** — [batch, traj_set, samples, waypoints, xyz]'},
                {'text': '· pred_rot: **(B, 1, 6, 64, 3, 3)** — SO(3) 회전 행렬'},
                {'text': '· **값의 의미:** 6개 미래 경로 후보 (미터 단위, ego-centric 좌표)'},
            ],
        },
        {
            'label': '',
            'lines': [
                {'text': '**평가 메트릭:** minADE = min_k (1/64 × Σ ||pred - gt||₂)'},
                {'text': '· 6개 중 ground truth에 가장 가까운 trajectory 선택'},
                {'text': '· 자율주행의 다중 모달 미래 분포를 커버리지 관점에서 평가'},
            ],
            'after_space': False,
        },
    ]
    add_bold_label_text(phs[1], sections, font_size=14)

    # ── 슬라이드 12: 파라미터 & 메모리 요약 ─────────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "파라미터 & 메모리 요약"

    ph1 = phs[1]
    left = ph1.left
    top = ph1.top
    width = ph1.width

    table_height = Emu(int(ph1.height * 0.58))

    rows, cols = 7, 4
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, table_height)
    table = table_shape.table

    col_widths = [Emu(3500000), Emu(2500000), Emu(2000000), Emu(3525958)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ['모듈', '파라미터', '비중', 'BF16 크기']
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h

    data = [
        ['Vision Encoder', '~575M', '5.2%', '1.15 GB'],
        ['VLM (Qwen3-VL-8B)', '~7,585M', '68.5%', '15.17 GB'],
        ['VLM LM Head', '~640M', '5.8%', '1.28 GB'],
        ['Expert Decoder', '~2,280M', '20.6%', '4.56 GB'],
        ['Diffusion Head', '~50M', '0.5%', '0.10 GB'],
        ['합계', '~11,130M', '100%', '22.16 GB'],
    ]
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val

    style_table_header(table, cols, font_size=12)
    style_table_data_rows(table, rows, cols, font_size=12)

    # 테이블 아래 추가 텍스트 (TextBox)
    text_top = top + table_height + Emu(250000)
    text_height = Emu(1000000)
    txBox = slide.shapes.add_textbox(left, text_top, width, text_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = '· 모델 파라미터 합계: 22.16 GB  |  Peak VRAM (추론 시): 21.52 GB'
    r1.font.size = Pt(13)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = '· VLM + LM Head가 전체의 74.3% (16.45 GB) → 12GB VRAM 초과, 메모리 최적화 핵심 대상'
    r2.font.size = Pt(13)
    r2.font.bold = True

    # placeholder 제거
    sp = ph1._element
    sp.getparent().remove(sp)

    # ── 슬라이드 13: 텐서 흐름 종합 다이어그램 ──────────────────────
    slide = prs.slides.add_slide(layouts['기본'])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    phs[0].text = "텐서 흐름 종합 다이어그램"

    flow_lines = [
        'Image (H×W×3)',
        '  → Patch Embed → (N, 1152)',
        '    → 27× ViT → (N, 1152)',
        '      → Merger → (N/4, 4096)',
        '        → [+ Text + History] → (B, L, 4096)',
        '          → 36× VLM → KV cache (36L)',
        '            → action_in_proj → (B, 64, 2048)',
        '              → 16× Expert (×10 steps) → (B, 64, 2048)',
        '                → action_out_proj → (B, 64, 2)',
        '                  → Unicycle → (B, 64, 3) [x, y, yaw]',
        '                    → 6× sampling → (B, 1, 6, 64, 3)',
    ]

    tf = phs[1].text_frame
    tf.word_wrap = False
    tf.clear()

    for i, line in enumerate(flow_lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.name = 'Consolas'

    # 저장
    prs.save(OUTPUT_PATH)
    print(f"저장 완료: {OUTPUT_PATH}")
    print(f"슬라이드 수: {len(prs.slides)}")


if __name__ == '__main__':
    main()
