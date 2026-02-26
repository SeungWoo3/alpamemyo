#!/usr/bin/env python3
"""
meeting-notes-01.md 내용을 AVEES Template에 맞춰 PPTX로 생성
"""
import os
import sys
import zipfile
import shutil
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ── 경로 설정 ──
TEMPLATE_PATH = "/mnt/c/Users/seungwoo/Desktop/AVEES Template_0.1.potx"
OUTPUT_PATH = "/home/seungwoo/workspace/docs/seminar/meeting-notes-01.pptx"
FIGURES_DIR = "/home/seungwoo/workspace/docs/seminar/figures"

# ── 색상 팔레트 ──
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x00, 0x70, 0xC0)  # 파란색 강조
RED = RGBColor(0xC0, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
TABLE_HEADER_BG = RGBColor(0x2D, 0x5F, 0x8A)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)


def fix_potx_to_pptx(potx_path):
    """POTX 파일의 Content_Types를 수정하여 PPTX로 사용 가능하게 만듦"""
    tmp = tempfile.mktemp(suffix=".pptx")
    shutil.copy2(potx_path, tmp)

    with zipfile.ZipFile(tmp, "r") as zin:
        content_types = zin.read("[Content_Types].xml").decode("utf-8")

    content_types = content_types.replace(
        "presentationml.template.main+xml",
        "presentationml.presentation.main+xml"
    )

    with zipfile.ZipFile(tmp, "r") as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}

    data["[Content_Types].xml"] = content_types.encode("utf-8")

    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for name in names:
            zout.writestr(name, data[name])

    return tmp


def set_cell_text(cell, text, font_size=11, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    """테이블 셀에 텍스트 설정"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 마진
    cell.margin_left = Emu(72000)
    cell.margin_right = Emu(72000)
    cell.margin_top = Emu(36000)
    cell.margin_bottom = Emu(36000)


def set_cell_bg(cell, color):
    """셀 배경색 설정"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": f"{color}"})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def add_table(slide, rows_data, left, top, width, height, col_widths=None, font_size=11, header_font_size=12):
    """슬라이드에 표 추가"""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 열 너비 설정
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # 헤더 행
    for j, val in enumerate(rows_data[0]):
        cell = table.cell(0, j)
        set_cell_text(cell, val, font_size=header_font_size, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        set_cell_bg(cell, "2D5F8A")

    # 데이터 행
    for i in range(1, n_rows):
        for j, val in enumerate(rows_data[i]):
            cell = table.cell(i, j)
            is_bold = j == 0  # 첫 열 볼드
            alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            color = DARK
            # 강조 표시
            if isinstance(val, str) and val.startswith("**") and val.endswith("**"):
                val = val[2:-2]
                is_bold = True
                color = RED
            set_cell_text(cell, val, font_size=font_size, bold=is_bold, color=color, alignment=alignment)
            if i % 2 == 0:
                set_cell_bg(cell, "E8F0F8")

    return table_shape


def add_bullet_slide(slide, title_text, bullets, sub_bullets=None):
    """기본 레이아웃 슬라이드에 제목과 불릿 포인트 추가"""
    # 제목
    title = slide.placeholders[0]
    title.text = title_text

    # 본문
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)

        # 볼드 처리
        if bullet.startswith("**") and "**" in bullet[2:]:
            parts = bullet.split("**")
            for idx, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                run.text = part
                run.font.size = Pt(18)
                run.font.name = "맑은 고딕"
                run.font.bold = (idx % 2 == 1)
        else:
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(18)
            run.font.name = "맑은 고딕"

        # 서브 불릿
        if sub_bullets and i in sub_bullets:
            for sub in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.level = 1
                sp.space_after = Pt(4)
                sr = sp.add_run()
                sr.text = sub
                sr.font.size = Pt(14)
                sr.font.name = "맑은 고딕"
                sr.font.color.rgb = GRAY


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    return txBox


def main():
    # 1) 템플릿 로드
    fixed_template = fix_potx_to_pptx(TEMPLATE_PATH)
    prs = Presentation(fixed_template)

    layout_cover = prs.slide_layouts[0]      # 표지
    layout_basic = prs.slide_layouts[1]      # 기본
    layout_section = prs.slide_layouts[4]    # 간지
    layout_blank = prs.slide_layouts[3]      # 공백
    layout_qa = prs.slide_layouts[5]         # Q&A

    # ═══════════════════════════════════════════
    # 슬라이드 1: 표지
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(layout_cover)
    slide.placeholders[0].text = "Alpamayo-R1\nVRAM 최적화 연구 보고서"
    slide.placeholders[1].text = "12GB GPU에서의 VLA 모델 추론을 위한\n파라미터 스왑 방법론 비교 분석"
    slide.placeholders[10].text = "AVEES Lab"
    if 2 in slide.placeholders:
        slide.placeholders[2].text = "2026-02-26"

    # ═══════════════════════════════════════════
    # 슬라이드 2: 목차
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "목차"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    toc_items = [
        ("1.", "Alpamayo-R1 VRAM 사용량 분석"),
        ("2.", "스왑 오버헤드 분석"),
        ("3.", "기존 오프로딩 프레임워크의 한계"),
        ("4.", "Demand Layering 적용 결과"),
        ("5.", "스왑 방법론 비교 (1~4)"),
        ("6.", "실행시간 비교 종합"),
        ("7.", "WSL2 환경 오버헤드"),
    ]
    for i, (num, title) in enumerate(toc_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.level = 0
        run_num = p.add_run()
        run_num.text = f"  {num}  "
        run_num.font.size = Pt(22)
        run_num.font.bold = True
        run_num.font.color.rgb = ACCENT
        run_num.font.name = "맑은 고딕"
        run_title = p.add_run()
        run_title.text = title
        run_title.font.size = Pt(22)
        run_title.font.name = "맑은 고딕"

    # ═══════════════════════════════════════════
    # 섹션 1: VRAM 사용량 분석
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(layout_section)
    slide.placeholders[0].text = "1. VRAM 사용량 분석"

    # 슬라이드: 모델 구조
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "Alpamayo-R1 모델 구조"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "11.08B 파라미터, BF16 기준 ~22.16 GB의 VLA 모델"
    run.font.size = Pt(16)
    run.font.name = "맑은 고딕"
    run.font.color.rgb = GRAY

    table_data = [
        ["모듈", "레이어 수", "크기 (BF16)", "비율"],
        ["Vision Encoder", "—", "1.15 GB", "5.2%"],
        ["VLM (Qwen3-VL-8B)", "36", "15.17 GB", "68.5%"],
        ["VLM LM Head", "—", "1.28 GB", "5.8%"],
        ["Expert (Trajectory Decoder)", "36", "4.56 GB", "20.6%"],
        ["합계", "—", "~22.16 GB", "100%"],
    ]
    add_table(slide, table_data,
              left=Inches(0.8), top=Inches(2.2), width=Inches(11.5), height=Inches(3.5),
              col_widths=[Inches(4.5), Inches(2), Inches(2.5), Inches(2.5)])

    # 슬라이드: VRAM 점유 비율
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "추론 시 VRAM 점유 비율"
    # 왼쪽에 테이블
    table_data = [
        ["분류", "구성요소", "크기", "비율"],
        ["정적 (파라미터)", "VLM Language Model", "15.17 GB", "63.7%"],
        ["", "Vision Encoder", "1.15 GB", "4.8%"],
        ["", "VLM LM Head", "1.28 GB", "5.4%"],
        ["", "Expert (Decoder)", "4.56 GB", "19.1%"],
        ["동적 (추론 중)", "KV Cache", "0.56 GB", "2.4%"],
        ["", "Activation + Overhead", "1.10 GB", "4.6%"],
    ]
    add_table(slide, table_data,
              left=Inches(0.4), top=Inches(1.2), width=Inches(8.0), height=Inches(4.2),
              col_widths=[Inches(2.2), Inches(2.8), Inches(1.5), Inches(1.5)], font_size=10, header_font_size=11)

    # 핵심 메시지 텍스트 박스
    add_text_box(slide, Inches(8.7), Inches(1.5), Inches(4.0), Inches(1.5),
                 "파라미터 93%\nvs\n동적 7%", font_size=28, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(8.7), Inches(3.3), Inches(4.0), Inches(1.0),
                 "모델 파라미터가 VRAM의\n대부분을 차지", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.4), Inches(5.5), Inches(12.0), Inches(1.0),
                 "측정 Peak: 21.52 GB  |  이론 합계: ~23.82 GB  |  정적 93.0% vs 동적 7.0%",
                 font_size=13, color=GRAY, alignment=PP_ALIGN.LEFT)

    # 슬라이드: KV Cache 비교
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "일반 LLM vs Alpamayo: KV Cache 비중"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    table_data = [
        ["", "일반 LLM (예: Llama-70B)", "Alpamayo-R1"],
        ["KV Cache 비중", "30-50%+", "2.4%"],
        ["원인", "긴 컨텍스트 + Full MHA", "GQA + 짧은 CoT (~15토큰)"],
        ["KV Cache 크기", "수~수십 GB", "0.56 GB"],
        ["메모리 병목", "KV Cache + 파라미터", "파라미터 단독 (93%)"],
    ]
    add_table(slide, table_data,
              left=Inches(0.5), top=Inches(1.2), width=Inches(12.0), height=Inches(3.2),
              col_widths=[Inches(3.0), Inches(4.5), Inches(4.5)], font_size=13, header_font_size=14)

    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(12.0), Inches(1.5),
                 "핵심: Alpamayo에서는 KV Cache 최적화가 아닌, 파라미터 자체의 메모리 최적화가 핵심 과제\n"
                 "→ VLM 단독 15.17GB > 12GB VRAM이므로 레이어 단위 파라미터 스와핑 필수",
                 font_size=15, bold=True, color=ACCENT)

    # ═══════════════════════════════════════════
    # 섹션 2: 스왑 오버헤드 분석
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(layout_section)
    slide.placeholders[0].text = "2. 스왑 오버헤드 분석"

    # 슬라이드: Vanilla Baseline
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "Vanilla Alpamayo (BF16, Unified Memory)"
    add_bullet_slide(slide, "Vanilla Alpamayo (BF16, Unified Memory)", [
        "환경: RTX 3080 Ti 12GB, WSL2, PyTorch 2.8.0+cu128",
        "**실행 결과: 273.79s** (이론 최적 ~5s 대비 **55배 느림**)",
        "모델 ~22GB > 12GB VRAM → Unified Memory로 CPU RAM 사용",
        "4KB~2MB 단위 on-demand 페이지 폴트 → PCIe 효율 저하",
        "VRAM 대역폭(912 GB/s) vs PCIe Gen3(8.5 GB/s) = **107배 격차**",
    ])

    # 슬라이드: 이론 최적값
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "이론 최적 추론 시간 (스왑 없는 경우)"
    table_data = [
        ["단계", "읽기량", "반복", "시간"],
        ["VLM (토큰 생성)", "16.45 GB/토큰", "256회", "4.62s"],
        ["Expert (디퓨전)", "4.56 GB/스텝", "10회", "0.05s"],
        ["Vision + KV Cache 등", "", "", "~0.3s"],
        ["합계", "", "", "≈ 5.0s"],
    ]
    add_table(slide, table_data,
              left=Inches(0.8), top=Inches(1.3), width=Inches(10.5), height=Inches(2.8),
              col_widths=[Inches(3.5), Inches(2.5), Inches(2.0), Inches(2.5)])

    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(10.5), Inches(2.0),
                 "연산 시간(0.44ms/토큰)은 읽기 시간(18ms/토큰)의 2.4%로 무시 가능\n"
                 "→ Memory-bandwidth-bound 워크로드\n\n"
                 "VRAM 내부 912 GB/s vs PCIe Gen3 8.5 GB/s = 107배 격차",
                 font_size=15, color=GRAY)

    # 슬라이드: 4-bit VRAM 제한 실험
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "4-bit 양자화 + VRAM 제한 실험"
    table_data = [
        ["VRAM 제한", "가용 VRAM", "추론 시간", "Peak VRAM", "비고"],
        ["12 GB", "12.00 GB", "4.79s", "8.87 GB", "Baseline"],
        ["11 GB", "11.30 GB", "5.29s", "9.59 GB", ""],
        ["10 GB", "10.29 GB", "6.96s", "10.59 GB", ""],
        ["9 GB", "9.29 GB", "111.36s", "11.58 GB", "Cliff"],
        ["8 GB", "8.29 GB", "179.46s", "12.59 GB", ""],
        ["7 GB", "7.30 GB", "373.05s", "13.61 GB", "재측정"],
        ["6 GB", "6.29 GB", "1506.88s", "14.59 GB", ""],
    ]
    add_table(slide, table_data,
              left=Inches(0.5), top=Inches(1.2), width=Inches(12.0), height=Inches(5.0),
              col_widths=[Inches(1.8), Inches(2.0), Inches(2.5), Inches(2.5), Inches(3.2)],
              font_size=12, header_font_size=13)

    # 슬라이드: Performance Cliff 시각화
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "Performance Cliff 분석"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    cliff_img = os.path.join(FIGURES_DIR, "performance_cliff.png")
    if os.path.exists(cliff_img):
        slide.shapes.add_picture(cliff_img, Inches(0.5), Inches(1.2), width=Inches(7.0))

    add_text_box(slide, Inches(7.8), Inches(1.5), Inches(4.8), Inches(3.5),
                 "10GB → 9GB 구간에서\n16배 급감\n(6.96s → 111.36s)\n\n"
                 "가용 VRAM이 모델 크기\n(8.87GB)를 밑돌기 시작하면\n"
                 "Unified Memory 스왑 발생\n→ 불연속적 성능 급락",
                 font_size=15, color=DARK)

    # ═══════════════════════════════════════════
    # 섹션 3: 기존 오프로딩 프레임워크의 한계
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(layout_section)
    slide.placeholders[0].text = "3. 기존 오프로딩 프레임워크의 한계"

    # 슬라이드: device_map="auto"
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "device_map=\"auto\" (HuggingFace Accelerate)"
    add_bullet_slide(slide, "device_map=\"auto\" (HuggingFace Accelerate)", [
        "VRAM 부족 시 표준 해결책: **device_map=\"auto\"**",
        "VRAM에 들어가는 레이어는 GPU, 초과분은 CPU/Disk에 배치",
        "**실험 결과: 추론 실패** (CUDA device-side assert)",
        "**원인**: fuse_traj_tokens()의 masked_scatter가 디바이스 분산과 충돌",
    ], sub_bullets={
        2: ["VLM 레이어 0-16 → GPU (4.49B, 40.6%)",
            "VLM 레이어 17-35 + LM Head + Expert → CPU (6.58B, 59.4%)"],
    })

    # 슬라이드: 커스텀 구현 필요성
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "커스텀 오프로딩 구현이 필요한 이유"

    table_data = [
        ["문제", "설명"],
        ["fuse_traj_tokens() 비호환", "masked_scatter가 디바이스 분산과 충돌"],
        ["KV Cache 공유", "Expert가 VLM의 past_key_values를 직접 참조"],
        ["generate() 내부 가정", "모든 모듈이 동일 디바이스에 있다고 가정"],
    ]
    add_table(slide, table_data,
              left=Inches(0.5), top=Inches(1.2), width=Inches(12.0), height=Inches(2.8),
              col_widths=[Inches(4.0), Inches(8.0)], font_size=13, header_font_size=14)

    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12.0), Inches(2.0),
                 "결론: 기존 프레임워크(device_map, Accelerate)의 자동 오프로딩은\n"
                 "Alpamayo에서 동작하지 않음 → 레이어 단위 커스텀 오프로딩 구현 필수",
                 font_size=16, bold=True, color=ACCENT)

    # device_map 분석 시각화
    dm_img = os.path.join(FIGURES_DIR, "02_device_map_analysis.png")
    if os.path.exists(dm_img):
        slide.shapes.add_picture(dm_img, Inches(8.5), Inches(1.2), width=Inches(4.0))

    # ═══════════════════════════════════════════
    # 섹션 4: Demand Layering
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(layout_section)
    slide.placeholders[0].text = "4. Demand Layering 적용 결과"

    # 슬라이드: 논문 소개
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "Demand Layering 논문 (RTSS 2022)"
    add_bullet_slide(slide, "Demand Layering 논문 (RTSS 2022)", [
        "DNN 추론 시 모델 파라미터를 레이어 단위로 로딩/실행",
        "**3-Stage Pipeline**: Read(SSD→staging) → Copy(staging→GPU) → Kernel",
        "원 논문: iGPU + NVMe SSD 환경 (통합 메모리)",
        "본 연구: dGPU + CPU RAM → Sequential (H2D→Execute→Free)",
        "D2H 불필요: CPU에 원본 보유 → GPU free만 수행",
    ])

    # 슬라이드: 실험 결과
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "Demand Layering 실험 결과"
    table_data = [
        ["항목", "값"],
        ["추론 시간", "43.38s"],
        ["Peak VRAM", "11.03 GB"],
        ["GPU 상주", "9.86 GB"],
        ["H2D 전송", "450회, 총 24.23s (avg 53.83 ms)"],
        ["GPU Free", "450회, 총 1.05s (avg 2.32 ms)"],
        ["vs Baseline (273.79s)", "6.31x 빠름"],
    ]
    add_table(slide, table_data,
              left=Inches(0.5), top=Inches(1.2), width=Inches(6.5), height=Inches(4.5),
              col_widths=[Inches(3.0), Inches(3.5)], font_size=13, header_font_size=14)

    # VLM 구성 설명
    add_text_box(slide, Inches(7.3), Inches(1.5), Inches(5.2), Inches(3.0),
                 "구현 방식:\n"
                 "- VLM 36개 레이어 중 30개 오프로드\n"
                 "- 6개는 GPU 상주\n"
                 "- forward hook: H2D → Execute → Free\n"
                 "- D2H 제거: CPU 원본 참조 복원\n\n"
                 "VLM 30 레이어 × 15 passes\n= 450 H2D + 450 Free",
                 font_size=14, color=GRAY)

    # 슬라이드: 한계 분석
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "Demand Layering 한계 분석"

    table_data = [
        ["구성 요소", "시간", "비율"],
        ["H2D 전송", "24.23s", "55.9%"],
        ["계산", "~18.1s", "41.7%"],
        ["GPU Free", "1.05s", "2.4%"],
        ["합계", "43.38s", "100%"],
    ]
    add_table(slide, table_data,
              left=Inches(0.5), top=Inches(1.2), width=Inches(5.5), height=Inches(3.0),
              col_widths=[Inches(2.0), Inches(1.5), Inches(2.0)], font_size=13, header_font_size=14)

    add_text_box(slide, Inches(6.3), Inches(1.3), Inches(6.0), Inches(4.5),
                 "핵심 관찰:\n\n"
                 "1. H2D가 추론 시간의 56% — PCIe 병목\n\n"
                 "2. GPU Free는 2.4%로 무시 가능\n"
                 "   (avg 2.32ms, D2H 대비 67배 빠름)\n\n"
                 "3. 동기식 전송 → 비동기 오버랩으로\n"
                 "   H2D 24.2s를 계산(18.1s)에 은닉 가능\n\n"
                 "4. Per-layer: H2D 54ms + Compute 40ms\n"
                 "   + Free 2ms = ~96ms",
                 font_size=15, color=DARK)

    # ═══════════════════════════════════════════
    # 섹션 5: 스왑 방법론 비교
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(layout_section)
    slide.placeholders[0].text = "5. 스왑 방법론 비교"

    # 슬라이드: 방법 1
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "방법 1: 레이어 단위 동기 스왑"
    add_bullet_slide(slide, "방법 1: 레이어 단위 동기 스왑", [
        "개별 Transformer 레이어를 forward hook으로 H2D → Execute → Free",
        "레이어당: H2D 54ms → Compute 40ms → Free 2ms (총 ~96ms)",
        "**실측: 43.38s** (vs Baseline 273.79s → **6.31x 가속**)",
        "D2H 불필요: CPU 원본 참조 복원 + GPU 메모리 해제",
        "구현 완료: research/07-demand-layering-no-d2h/",
    ])

    # 슬라이드: 방법 2
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "방법 2: 최적 청크 전송 (2-8MB)"
    table_data = [
        ["청크 크기", "유효 대역폭", "비고"],
        ["2-8 MB", "12.16-12.41 GB/s", "최적 구간"],
        ["16 MB", "12.08 GB/s", "소폭 감소"],
        ["512 MB (단일)", "7.77 GB/s", "기준"],
        ["PCIe Gen3 이론", "15.75 GB/s", "실측 79% 활용"],
    ]
    add_table(slide, table_data,
              left=Inches(0.5), top=Inches(1.2), width=Inches(8.0), height=Inches(3.0),
              col_widths=[Inches(2.5), Inches(3.0), Inches(2.5)], font_size=13, header_font_size=14)

    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12.0), Inches(2.0),
                 "Pinned memory + 2-8MB 청크 → 단일 전송 대비 1.6x 대역폭 향상\n"
                 "이론 시간: 24.23s / 1.6 ≈ 15s + compute 18.1s + free 1.0s ≈ ~34s",
                 font_size=16, bold=True, color=ACCENT)

    # 슬라이드: 방법 3
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "5.3 방법 3: 모듈 단위 스왑 → 불가능"
    add_bullet_slide(slide, "5.3 방법 3: 모듈 단위 스왑 → 불가능", [
        "Vision / VLM / Expert를 모듈 통째로 스왑하는 방식",
        "**VLM 단독 15.17GB > 12GB VRAM** — 모듈 통째 로드 불가",
        "Vision Encoder (1.15GB) → 가능",
        "Expert (4.56GB) → 가능",
    ], sub_bullets={
        1: ["순수 모듈 단위 스왑은 구조적으로 불가능"],
    })

    # 슬라이드: 교수님 아이디어 (4.4)
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "5.4 교수님 아이디어: 2-Stage 스와핑"
    add_bullet_slide(slide, "5.4 교수님 아이디어: 2-Stage 스와핑", [
        "**A. 모듈 단위 2-Stage → 불가능**",
        "**B. 레이어 단위 2-Stage → 유효**",
        "**결론**: 방법 1의 현재 구현이 이미 레이어 단위 2-Stage에 해당",
    ], sub_bullets={
        0: ["Stage 1: Vision+VLM → GPU / Stage 2: Expert → GPU",
            "VLM 단독 15.17GB > 12GB → 불가능"],
        1: ["Stage 1: VLM 레이어별 스왑 (prefill + token generation)",
            "Stage 2: Expert 전체 GPU 로드 (4.56GB, 적재 가능) → diffusion 실행"],
    })

    # 슬라이드: 방법 4 (4.5)
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "5.5 방법 4: 비동기 2-Stream 파이프라인"
    add_bullet_slide(slide, "5.5 방법 4: 비동기 2-Stream 파이프라인", [
        "방법 1의 H2D(54ms)와 Compute(40ms)를 비동기 오버랩",
        "2개 CUDA 스트림으로 계산과 H2D prefetch를 겹침",
        "Compute Stream: 현재 레이어 forward pass",
        "H2D Stream: 다음 레이어 prefetch + 이전 레이어 free",
    ])

    # 슬라이드: 방법 4 이론 시간
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "방법 4: 이론적 성능 분석"
    table_data = [
        ["시나리오", "Per-layer", "450 layers", "총 시간"],
        ["기본 비동기 (Pageable)", "max(54ms, 40ms) = 54ms", "450 × 54ms", "~24s"],
        ["+ Chunked (12.4 GB/s)", "max(31ms, 40ms) = 40ms", "450 × 40ms", "~18s"],
    ]
    add_table(slide, table_data,
              left=Inches(0.5), top=Inches(1.2), width=Inches(12.0), height=Inches(2.2),
              col_widths=[Inches(3.5), Inches(3.5), Inches(2.5), Inches(2.5)], font_size=13, header_font_size=14)

    # 선행 연구 비교
    table_data2 = [
        ["시스템", "연도", "D2H 사용 이유", "대상"],
        ["SwapAdvisor (ASPLOS)", "2020", "Activation 저장 (학습)", "학습"],
        ["Demand Layering (RTSS)", "2022", "SSD writeback", "RT 추론"],
        ["FlexGen (ICML)", "2023", "KV Cache + Activation", "배치 추론"],
        ["PIPO", "2025", "KV Cache saving", "소비자 추론"],
        ["본 연구", "2026", "불필요 (free만)", "VLA 추론"],
    ]
    add_table(slide, table_data2,
              left=Inches(0.5), top=Inches(3.8), width=Inches(12.0), height=Inches(3.0),
              col_widths=[Inches(3.5), Inches(1.5), Inches(4.0), Inches(3.0)], font_size=11, header_font_size=12)

    # 슬라이드: 방법 4 실험 결과
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "방법 4: 실험 결과"
    table_data = [
        ["Variant", "추론 시간", "Peak VRAM", "핵심 관찰"],
        ["A: Pageable async", "44.16s (1.0x)", "11.15 GB", "non_blocking 무효 (CUDA 내부 동기)"],
        ["C: Staged pinned buffer", "38.45s (1.13x)", "11.38 GB", "DMA 오버랩 성공, CPU staging 병목"],
    ]
    add_table(slide, table_data,
              left=Inches(0.5), top=Inches(1.2), width=Inches(12.0), height=Inches(2.2),
              col_widths=[Inches(3.0), Inches(2.5), Inches(2.0), Inches(4.5)], font_size=13, header_font_size=14)

    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(12.0), Inches(3.0),
                 "핵심 발견:\n"
                 "• DMA 오버랩 성공: wait_dma avg 0.05ms (DMA가 compute 중 완료)\n"
                 "• 병목 이동: H2D DMA → CPU memcpy (pageable→pinned, avg 44ms)\n"
                 "• Pageable에서 non_blocking=True는 CUDA가 내부 bounce buffer로 동기 처리\n"
                 "• 전체 파라미터 pinning(11.6GB) 불가: 16GB RAM 한계\n"
                 "• 최적 조건(충분한 RAM + pinned memory) 시 이론 ~24s 가능",
                 font_size=14, color=DARK)

    # ═══════════════════════════════════════════
    # 섹션 6: 종합 비교
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(layout_section)
    slide.placeholders[0].text = "6. 이론적 실행시간 비교 종합"

    # 슬라이드: 종합 테이블
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "실행시간 비교 종합"
    table_data = [
        ["방법", "추론 시간", "vs Baseline", "비고"],
        ["Baseline (BF16 Unified Memory)", "273.79s", "1x", "실측"],
        ["방법 1 (동기 H2D + Free)", "43.38s", "6.31x", "실측"],
        ["방법 2 (청크 전송)", "~34s", "8.1x", ""],
        ["방법 3 (모듈 단위)", "N/A", "—", "VLM>12GB, 불가"],
        ["방법 4 (비동기, staged pinned)", "38.45s", "7.1x", "실측"],
        ["방법 4 (이론, 전체 pinned)", "~24s", "11.4x", "이론"],
        ["방법 4 + 2 (비동기 + 청크)", "~18s", "15.2x", "이론 최선"],
        ["이론 최적 (스왑 없음)", "~5s", "55x", "24GB+ GPU"],
    ]
    add_table(slide, table_data,
              left=Inches(0.3), top=Inches(1.1), width=Inches(12.5), height=Inches(5.2),
              col_widths=[Inches(4.5), Inches(2.5), Inches(2.5), Inches(3.0)],
              font_size=13, header_font_size=14)

    # 슬라이드: 핵심 메시지
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "현실적 목표와 한계"
    add_bullet_slide(slide, "현실적 목표와 한계", [
        "**현실적 목표**: 방법 4+2로 273.79s → ~18s (**15배 가속**)",
        "이론 최적(~5s) 대비 3.6배 차이는 PCIe Gen3의 구조적 한계",
        "**격차 해소 방법**:",
    ], sub_bullets={
        2: [
            "VRAM 증설 (24GB+ GPU)",
            "더 빠른 인터커넥트 (PCIe Gen5: 2배, NVLink: 10배+)",
            "모델 축소 — 양자화 (연구에서 배제)",
        ]
    })

    # ═══════════════════════════════════════════
    # 섹션 7: WSL2 오버헤드
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(layout_basic)
    slide.placeholders[0].text = "WSL2 환경 오버헤드"
    table_data = [
        ["항목", "WSL2 (현재)", "네이티브 Linux (예상)", "차이"],
        ["Pageable D2H 대역폭", "2.48 GB/s", "~11 GB/s", "4.5배 느림"],
        ["Pinned D2H 대역폭", "11.8 GB/s", "~12+ GB/s", "유사"],
        ["Pageable H2D 대역폭", "7.77 GB/s", "~8+ GB/s", "유사"],
        ["per-transfer 오버헤드", "0.36 ms", "~0.05 ms", "7배"],
        ["cudaMemPrefetchAsync(CPU)", "미지원", "지원", "—"],
    ]
    add_table(slide, table_data,
              left=Inches(0.4), top=Inches(1.2), width=Inches(12.2), height=Inches(3.5),
              col_widths=[Inches(3.5), Inches(2.7), Inches(3.0), Inches(3.0)], font_size=12, header_font_size=13)

    add_text_box(slide, Inches(0.4), Inches(5.0), Inches(12.0), Inches(1.5),
                 "방법 4는 Pinned memory 사용 → WSL2 오버헤드 영향 최소\n"
                 "최적화 방향성 탐색에는 WSL2에서도 유효, 최종 평가는 네이티브 Linux 권장",
                 font_size=14, color=GRAY)

    # ═══════════════════════════════════════════
    # Q&A
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(layout_qa)

    # ── 저장 ──
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"PPTX 생성 완료: {OUTPUT_PATH}")
    print(f"파일 크기: {os.path.getsize(OUTPUT_PATH) / 1024:.1f} KB")
    print(f"총 슬라이드: {len(prs.slides)}")

    # 임시 파일 정리
    os.unlink(fixed_template)


if __name__ == "__main__":
    main()
